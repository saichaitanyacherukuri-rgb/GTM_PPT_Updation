"""
PPT Auto-Updater: Reads CSV/Excel data files and updates corresponding
tables in an existing PowerPoint presentation, preserving all formatting.
"""

import json
import os
import re
import sys
from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from pptx.util import Pt


SCRIPT_DIR = Path(__file__).resolve().parent
CONFIG_FILE = SCRIPT_DIR / "slide_config.json"
MASTER_FILE_NAME = "GTM_GAI_Tabls.xlsx"

DEFAULTS = {
    "table_index": 0,
    "start_row": 1,
    "start_col": 0,
    "header_rows": 1,
    "skip_columns": [],
    "skip_rows": [],
    "sheet_name": None,
    "transpose": False,
    "font_name": "Avenir Next LT Pro",
    "font_size": 8,
    "cond_format_start_col": None,
    "cond_format_end_col": None,
    "cond_format_start_row": None,
}


def find_pptx(folder: Path, target_name: str = None) -> Path:
    """Find the .pptx file to update. If target_name is given, use that file."""
    if target_name:
        target = folder / target_name
        if not target.exists():
            sys.exit(f"ERROR: File not found: {target}")
        return target

    pptx_files = [
        f for f in folder.glob("*.pptx")
        if not f.name.startswith("~$")
        and "_updated" not in f.name
    ]
    if len(pptx_files) == 0:
        sys.exit("ERROR: No .pptx file found in the folder.")
    if len(pptx_files) > 1:
        names = ", ".join(f.name for f in pptx_files)
        sys.exit(f"ERROR: Multiple .pptx files found ({names}). "
                 f"Specify which one: py update_ppt.py \"filename.pptx\"")
    return pptx_files[0]


def find_data_files(folder: Path) -> dict[int, Path]:
    """
    Scan folder for files matching 'slide <number>.csv' or 'slide <number>.xlsx'.
    Returns a dict mapping slide number -> file path.
    """
    pattern = re.compile(r"^slide\s+(\d+)\.(csv|xlsx)$", re.IGNORECASE)
    result: dict[int, Path] = {}
    for f in folder.iterdir():
        if f.is_file():
            m = pattern.match(f.name)
            if m:
                slide_num = int(m.group(1))
                result[slide_num] = f
    return result


def parse_master_file(master_path: Path) -> dict:
    """Parse GTM_GAI_Tabls.xlsx and return a slide-sources dict.

    Sheet naming convention:
      "slide 10"    -> slide 10, single table
      "slide 15.1"  -> slide 15, first table
      "slide 15.2"  -> slide 15, second table

    Returns:
        dict[int, list[str]]  e.g. {10: ["slide 10"], 15: ["slide 15.1", "slide 15.2"]}
    """
    if not master_path.exists():
        sys.exit(f"ERROR: Master file not found: {master_path.name}\n"
                 f"Expected file: {master_path}")

    xls = pd.ExcelFile(master_path)
    sheet_pattern = re.compile(
        r"^slide\s+(\d+)(?:\.(\d+))?$", re.IGNORECASE
    )

    # Collect: slide_num -> list of (sub_index, sheet_name)
    raw: dict[int, list[tuple[int, str]]] = {}
    for sheet in xls.sheet_names:
        m = sheet_pattern.match(sheet.strip())
        if not m:
            continue
        slide_num = int(m.group(1))
        sub_idx   = int(m.group(2)) if m.group(2) else 0
        raw.setdefault(slide_num, []).append((sub_idx, sheet))

    if not raw:
        sys.exit(
            f"ERROR: No sheets matching 'slide N' or 'slide N.M' found in {master_path.name}.\n"
            f"Available sheets: {xls.sheet_names}"
        )

    # Sort each slide's sheets by sub-index
    result: dict[int, list[str]] = {}
    for slide_num, entries in raw.items():
        entries.sort(key=lambda x: x[0])
        result[slide_num] = [sheet for _, sheet in entries]

    return result


def load_config(config_file: Path = None) -> dict:
    """Load slide_config.json. Returns empty dict if not found."""
    cf = config_file or CONFIG_FILE
    if not cf.exists():
        print(f"WARNING: Config file not found at {cf}. Using defaults for all slides.")
        return {}
    with open(cf, "r", encoding="utf-8") as fh:
        raw = json.load(fh)
    # Strip internal documentation keys
    return {k: v for k, v in raw.items() if not k.startswith("_")}


def get_slide_config(config: dict, slide_num: int) -> dict:
    """Get merged config for a slide: user overrides on top of defaults."""
    user_cfg = config.get(str(slide_num), {})
    merged = {**DEFAULTS, **user_cfg}
    return merged


def get_sheet_names(file_path: Path) -> list[str]:
    """Return the list of sheet names for an Excel file, or ['Sheet1'] for CSV."""
    ext = file_path.suffix.lower()
    if ext == ".csv":
        return ["Sheet1"]
    elif ext == ".xlsx":
        xls = pd.ExcelFile(file_path)
        return xls.sheet_names
    return []


def read_data(file_path: Path, cfg: dict, sheet_name=None) -> pd.DataFrame:
    """Read a CSV or Excel file (specific sheet) into a DataFrame."""
    ext = file_path.suffix.lower()
    header_rows = cfg["header_rows"]
    header = list(range(header_rows)) if header_rows > 0 else None

    if ext == ".csv":
        df = pd.read_csv(file_path, header=header)
    elif ext == ".xlsx":
        sheet = sheet_name if sheet_name is not None else (cfg.get("sheet_name") or 0)
        df = pd.read_excel(file_path, header=header, sheet_name=sheet)
    else:
        sys.exit(f"ERROR: Unsupported file type: {ext}")

    if cfg.get("transpose"):
        df = df.T

    if isinstance(df.columns, pd.MultiIndex):
        df.columns = range(df.shape[1])

    df = df.reset_index(drop=True)
    df.columns = range(df.shape[1])
    return df


MONTH_ABBR = {
    1: "Jan", 2: "Feb", 3: "Mar", 4: "Apr", 5: "May", 6: "Jun",
    7: "Jul", 8: "Aug", 9: "Sep", 10: "Oct", 11: "Nov", 12: "Dec",
}


def format_month(value) -> str:
    """Convert a date/timestamp to 'Jan-25' format."""
    if pd.isna(value):
        return ""
    if isinstance(value, (pd.Timestamp, pd.DatetimeTZDtype)):
        return f"{MONTH_ABBR[value.month]}-{str(value.year)[2:]}"
    s = str(value).strip()
    try:
        dt = pd.to_datetime(s)
        return f"{MONTH_ABBR[dt.month]}-{str(dt.year)[2:]}"
    except Exception:
        return s


def format_usd_short(value) -> str:
    """Convert a number to short USD format: $234M, $2.34M, $234K, $45K."""
    if pd.isna(value):
        return ""
    try:
        num = float(value)
    except (ValueError, TypeError):
        return str(value)

    abs_num = abs(num)
    sign = "-" if num < 0 else ""

    if abs_num >= 1_000_000:
        m = abs_num / 1_000_000
        if m >= 100:
            return f"{sign}${m:.0f}M"
        elif m >= 10:
            return f"{sign}${m:.1f}M"
        else:
            return f"{sign}${m:.2f}M"
    elif abs_num >= 1_000:
        k = abs_num / 1_000
        if k >= 100:
            return f"{sign}${k:.0f}K"
        elif k >= 10:
            return f"{sign}${k:.1f}K"
        else:
            return f"{sign}${k:.2f}K"
    else:
        return f"{sign}${abs_num:.0f}"


def format_pct2(value) -> str:
    """Convert a decimal (0.2513) to percentage string with 2 decimals (25.13%)."""
    if pd.isna(value):
        return ""
    try:
        num = float(value)
    except (ValueError, TypeError):
        return str(value)
    return f"{num * 100:.2f}%"


def format_pct_auto(value) -> str:
    """Smart formatter: small decimals (abs < 10) become percentages, everything else passes through."""
    if pd.isna(value):
        return ""
    try:
        num = float(value)
    except (ValueError, TypeError):
        return str(value)
    if abs(num) < 10:
        return f"{num * 100:.1f}%"
    return f"{num:,.0f}"


def format_pct_auto_round(value) -> str:
    """Like pct_auto but with no decimals (rounded)."""
    import math
    if pd.isna(value):
        return ""
    try:
        num = float(value)
    except (ValueError, TypeError):
        return str(value)
    if abs(num) < 10:
        pct = round(num * 100, 10)  # eliminate float noise before rounding
        rounded = math.floor(pct + 0.5) if pct >= 0 else -math.floor(-pct + 0.5)
        return f"{int(rounded)}%"
    return f"{num:,.0f}"


def format_pct_halfup(value) -> str:
    """Like pct_auto_round but uses half-up rounding (>=0.5 always rounds up)."""
    import math
    if pd.isna(value):
        return ""
    try:
        num = float(value)
    except (ValueError, TypeError):
        return str(value)
    if abs(num) < 10:
        pct = round(num * 100, 10)  # eliminate float noise before rounding
        rounded = math.floor(pct + 0.5) if pct >= 0 else -math.floor(-pct + 0.5)
        return f"{int(rounded)}%"
    return f"{math.floor(abs(num) + 0.5) * (1 if num >= 0 else -1):,.0f}"


def format_usd_millions(value) -> str:
    """Always format as $X.XM (millions, 1 decimal). Never uses K."""
    if pd.isna(value):
        return ""
    try:
        num = float(str(value).replace(",", "").replace("$", ""))
    except (ValueError, TypeError):
        return str(value)
    sign = "-" if num < 0 else ""
    m = abs(num) / 1_000_000
    return f"{sign}${m:.1f}M"


def format_number_comma(value) -> str:
    """Format a number with comma thousands separator, no decimals.
    Passes through non-numeric values (e.g. already-formatted strings) unchanged."""
    if pd.isna(value):
        return ""
    try:
        num = float(str(value).replace(",", ""))
        return f"{num:,.0f}"
    except (ValueError, TypeError):
        return str(value)


def format_month_dm(value) -> str:
    """Format a date as DD-Mon (e.g. '24-Jan').
    Handles datetime/Timestamp objects and passes through already-formatted strings."""
    if pd.isna(value):
        return ""
    if isinstance(value, (pd.Timestamp,)):
        return value.strftime("%d-%b")
    try:
        from datetime import datetime as _dt
        if isinstance(value, _dt):
            return value.strftime("%d-%b")
    except Exception:
        pass
    try:
        return pd.to_datetime(value).strftime("%d-%b")
    except Exception:
        return str(value)


COLUMN_FORMATTERS = {
    "month": format_month,
    "month_dm": format_month_dm,
    "usd_short": format_usd_short,
    "usd_millions": format_usd_millions,
    "pct2": format_pct2,
    "pct_auto": format_pct_auto,
    "pct_auto_round": format_pct_auto_round,
    "pct_halfup": format_pct_halfup,
    "number_comma": format_number_comma,
}


def format_value(value, col_idx: int, col_formats: dict) -> str:
    """Apply column-specific formatting if configured, otherwise default to str."""
    fmt_name = col_formats.get(str(col_idx))
    if fmt_name and fmt_name in COLUMN_FORMATTERS:
        return COLUMN_FORMATTERS[fmt_name](value)
    if pd.isna(value):
        return ""
    return str(value)


def parse_numeric(text: str) -> float | None:
    """Try to extract a numeric value from text like '32.7%', '3,427', '280M'."""
    if not text or not text.strip():
        return None
    s = text.strip().replace(",", "").replace(" ", "")
    if s.endswith("%"):
        try:
            return float(s[:-1])
        except ValueError:
            return None
    for suffix, mult in [("M", 1e6), ("K", 1e3), ("B", 1e9)]:
        if s.endswith(suffix):
            try:
                return float(s[:-1]) * mult
            except ValueError:
                return None
    try:
        return float(s)
    except ValueError:
        return None


COLOR_LOW = (248, 105, 107)   # #F8696B — soft coral red
COLOR_MID = (255, 235, 132)   # #FFEB84 — warm yellow
COLOR_HIGH = (99, 190, 123)   # #63BE7B — muted green


def interpolate_color(val: float, min_val: float, max_val: float) -> RGBColor:
    """Red(low) -> Yellow(mid) -> Green(high) color scale using Excel-style palette."""
    if max_val == min_val:
        return RGBColor(*COLOR_MID)
    ratio = (val - min_val) / (max_val - min_val)
    ratio = max(0.0, min(1.0, ratio))
    if ratio < 0.5:
        t = ratio / 0.5
        r = int(COLOR_LOW[0] + (COLOR_MID[0] - COLOR_LOW[0]) * t)
        g = int(COLOR_LOW[1] + (COLOR_MID[1] - COLOR_LOW[1]) * t)
        b = int(COLOR_LOW[2] + (COLOR_MID[2] - COLOR_LOW[2]) * t)
    else:
        t = (ratio - 0.5) / 0.5
        r = int(COLOR_MID[0] + (COLOR_HIGH[0] - COLOR_MID[0]) * t)
        g = int(COLOR_MID[1] + (COLOR_HIGH[1] - COLOR_MID[1]) * t)
        b = int(COLOR_MID[2] + (COLOR_HIGH[2] - COLOR_MID[2]) * t)
    return RGBColor(r, g, b)


def clear_cell_fill(cell):
    """Remove all fill elements from a cell's tcPr and restore noFill so blank
    cells in the heat map range show no background color."""
    tc = cell._tc
    tc_pr = tc.find(qn("a:tcPr"))
    if tc_pr is None:
        return
    fill_tags = {
        qn("a:noFill"), qn("a:solidFill"), qn("a:gradFill"),
        qn("a:blipFill"), qn("a:pattFill"), qn("a:grpFill"),
    }
    for child in list(tc_pr):
        if child.tag in fill_tags:
            tc_pr.remove(child)
    # Insert noFill before cell3D/extLst to signal transparent background
    no_fill = tc_pr.makeelement(qn("a:noFill"), {})
    insert_before = {qn("a:cell3D"), qn("a:extLst")}
    for i, child in enumerate(tc_pr):
        if child.tag in insert_before:
            tc_pr.insert(i, no_fill)
            break
    else:
        tc_pr.append(no_fill)


def set_cell_fill(cell, color: RGBColor):
    """Set solid fill on a table cell.
    Removes ALL existing fill elements first (noFill, solidFill, gradFill, etc.)
    so there is never more than one fill element -- multiple fills are invalid OOXML.
    """
    tc = cell._tc
    tc_pr = tc.get_or_add_tcPr()
    fill_tags = {
        qn("a:noFill"), qn("a:solidFill"), qn("a:gradFill"),
        qn("a:blipFill"), qn("a:pattFill"), qn("a:grpFill"),
    }
    for child in list(tc_pr):
        if child.tag in fill_tags:
            tc_pr.remove(child)
    solid_fill = tc_pr.makeelement(qn("a:solidFill"), {})
    srgb = solid_fill.makeelement(qn("a:srgbClr"), {"val": str(color)})
    solid_fill.append(srgb)
    # Insert before cell3D and extLst to maintain correct OOXML element order
    insert_before = {qn("a:cell3D"), qn("a:extLst")}
    for i, child in enumerate(tc_pr):
        if child.tag in insert_before:
            tc_pr.insert(i, solid_fill)
            break
    else:
        tc_pr.append(solid_fill)


def _apply_cf_column(table, ppt_c: int, cf_start_row: int, ppt_rows: int) -> int:
    """Apply independent heat map to a single column. Returns number of colored cells."""
    col_values = []
    for ppt_r in range(cf_start_row, ppt_rows):
        val = parse_numeric(table.cell(ppt_r, ppt_c).text)
        if val is not None:
            col_values.append(val)

    if not col_values:
        return 0

    min_val = min(col_values)
    max_val = max(col_values)
    colored = 0
    for ppt_r in range(cf_start_row, ppt_rows):
        cell = table.cell(ppt_r, ppt_c)
        val = parse_numeric(cell.text)
        if val is None:
            clear_cell_fill(cell)
            continue
        set_cell_fill(cell, interpolate_color(val, min_val, max_val))
        colored += 1
    return colored


def apply_conditional_formatting(table, cfg: dict, start_row: int, ppt_rows: int, ppt_cols: int):
    """Apply red-yellow-green color scale to specified columns.

    Two modes (set via config):
    - Shared scale  (default): one min/max across all specified columns.
    - Per-column    (cond_format_per_column: true): each column uses its own min/max.
    """
    cf_start_col = cfg.get("cond_format_start_col")
    cf_end_col   = cfg.get("cond_format_end_col")
    cf_start_row = cfg.get("cond_format_start_row")
    per_column   = cfg.get("cond_format_per_column", False)

    if cf_start_col is None:
        return

    if cf_start_row is None:
        cf_start_row = start_row

    cf_end = (cf_end_col + 1) if cf_end_col is not None else ppt_cols
    col_range = range(cf_start_col, cf_end)

    # ── Per-column mode: each column has its own independent color scale ──
    if per_column:
        total_colored = 0
        for ppt_c in col_range:
            total_colored += _apply_cf_column(table, ppt_c, cf_start_row, ppt_rows)
        print(f"  Conditional formatting (per-column): applied to {total_colored} cells "
              f"(cols {cf_start_col}-{cf_end-1}, rows {cf_start_row}-{ppt_rows-1})")
        return

    # ── Shared scale mode: one min/max across all columns ──
    all_values = []
    for ppt_r in range(cf_start_row, ppt_rows):
        for ppt_c in col_range:
            val = parse_numeric(table.cell(ppt_r, ppt_c).text)
            if val is not None:
                all_values.append(val)

    if not all_values:
        return

    min_val = min(all_values)
    max_val = max(all_values)
    colored_count = 0

    for ppt_r in range(cf_start_row, ppt_rows):
        for ppt_c in col_range:
            cell = table.cell(ppt_r, ppt_c)
            val = parse_numeric(cell.text)
            if val is None:
                clear_cell_fill(cell)
                continue
            set_cell_fill(cell, interpolate_color(val, min_val, max_val))
            colored_count += 1

    print(f"  Conditional formatting: applied to {colored_count} cells "
          f"(cols {cf_start_col}-{cf_end-1}, rows {cf_start_row}-{ppt_rows-1}), "
          f"value range [{min_val:.1f} - {max_val:.1f}]")


def _xml_insert_ordered(parent, new_elem, *preceding_tags):
    """Insert new_elem before the first child whose tag is in preceding_tags.
    This enforces correct OOXML child-element ordering inside <a:rPr>.
    If none of the tags are found, the element is appended.
    """
    for i, child in enumerate(parent):
        if child.tag in preceding_tags:
            parent.insert(i, new_elem)
            return
    parent.append(new_elem)


def set_cell_text(cell, value: str, font_name: str = None, font_size: float = None,
                  bold: bool = None, font_color: RGBColor = None):
    """
    Update cell text while preserving formatting.
    Strategy: set the text on the first run, clear all other runs, then use
    direct XML to apply font properties to every <a:rPr> in the cell --
    this handles multi-run / line-break cells reliably.
    """
    from lxml import etree

    text_frame = cell.text_frame
    if not text_frame.paragraphs:
        text_frame.text = value
        return

    para = text_frame.paragraphs[0]

    if not para.runs:
        para.text = value

    if para.runs:
        para.runs[0].text = value
        for run in para.runs[1:]:
            run.text = ""

    for extra_para in text_frame.paragraphs[1:]:
        for run in extra_para.runs:
            run.text = ""
        if not extra_para.runs:
            extra_para.text = ""

    # Apply font properties to all runs in all paragraphs using python-pptx's
    # high-level API, which generates correct OOXML element ordering.
    for para in text_frame.paragraphs:
        for run in para.runs:
            if font_name:
                run.font.name = font_name
            if font_size:
                run.font.size = Pt(font_size)
            if bold is not None:
                run.font.bold = bold
            if font_color is not None:
                run.font.color.rgb = font_color


def get_tables_on_slide(slide) -> list:
    """Return all table shapes on a slide."""
    return [shape for shape in slide.shapes if shape.has_table]


def unmerge_table_cells(table, skip_cols: set = None, start_row: int = 0):
    """Remove horizontal and vertical merges so every cell is independent.
    Columns listed in skip_cols are left untouched (preserving their merges).
    Rows before start_row are left untouched (preserving header merges).
    """
    unmerged = 0
    for row_idx, row in enumerate(table.rows):
        if row_idx < start_row:
            continue
        for col_idx, cell in enumerate(row.cells):
            if skip_cols and col_idx in skip_cols:
                continue
            tc = cell._tc
            if tc.get("gridSpan") and tc.get("gridSpan") != "1":
                del tc.attrib["gridSpan"]
                unmerged += 1
            if tc.get("hMerge"):
                del tc.attrib["hMerge"]
                unmerged += 1
            if tc.get("rowSpan") and tc.get("rowSpan") != "1":
                del tc.attrib["rowSpan"]
                unmerged += 1
            if tc.get("vMerge"):
                del tc.attrib["vMerge"]
                unmerged += 1
    return unmerged


def update_slide_table(slide, slide_num: int, df: pd.DataFrame, cfg: dict):
    """Write DataFrame values into the specified table on the slide."""
    tables = get_tables_on_slide(slide)
    table_idx = cfg["table_index"]

    if not tables:
        print(f"  WARNING: Slide {slide_num} has no tables. Skipping.")
        return
    if table_idx >= len(tables):
        print(f"  WARNING: Slide {slide_num} has {len(tables)} table(s), "
              f"but config requests table_index={table_idx}. Skipping.")
        return

    table = tables[table_idx].table
    skip_cols = set(cfg["skip_columns"])
    merge_cols = set(cfg.get("merge_columns", []))
    no_unmerge = skip_cols | merge_cols
    start_row = cfg["start_row"]
    unmerged = unmerge_table_cells(table, skip_cols=no_unmerge, start_row=start_row)
    if unmerged:
        print(f"  Unmerged {unmerged} cell attributes to restore individual cells.")
    ppt_rows = len(table.rows)
    ppt_cols = len(table.columns)

    start_row = cfg["start_row"]
    start_col = cfg["start_col"]
    skip_rows = set(cfg["skip_rows"])
    font_name = cfg.get("font_name")
    font_size = cfg.get("font_size")
    col_formats = cfg.get("col_formats", {})
    negative_red = cfg.get("negative_red", False)
    bold_first_in_merge = cfg.get("bold_first_in_merge", False)
    always_bold_cols = set(cfg.get("always_bold_columns", []))
    bold_data_rows = set(cfg.get("bold_data_rows", []))
    # Column-specific font colouring (by PPT column index)
    pos_green_cols      = set(cfg.get("positive_green_columns", []))
    neg_red_cols        = set(cfg.get("negative_red_columns", []))
    positive_green      = cfg.get("positive_green", False)   # applies to ALL columns
    # Columns where colour is applied ONLY when the formatted value is a percentage
    pct_color_only_cols = set(cfg.get("pct_color_only_columns", []))
    # Row-level gate: only colour cells in rows where a data column matches a substring
    # e.g. {"col": 2, "contains": "Attainment"}
    color_row_filter    = cfg.get("color_only_rows_where", None)
    # Threshold-based coloring: raw_value >= threshold → green, < threshold → red
    # Takes priority over positive/negative logic when set
    color_threshold     = cfg.get("color_threshold", None)

    RED   = RGBColor(0xFF, 0x00, 0x00)
    GREEN = RGBColor(0x00, 0x70, 0x00)  # dark green for readability

    data_row_idx = 0
    for ppt_r in range(start_row, ppt_rows):
        if ppt_r in skip_rows:
            continue
        if data_row_idx >= len(df):
            break

        # Check row-level colour gate
        row_color_allowed = True
        if color_row_filter:
            filter_col  = color_row_filter.get("col", 0)
            filter_text = color_row_filter.get("contains", "").lower()
            if filter_col < df.shape[1]:
                cell_val = str(df.iloc[data_row_idx, filter_col]).lower()
                row_color_allowed = filter_text in cell_val
            else:
                row_color_allowed = False

        is_kpi_first = False
        if bold_first_in_merge and merge_cols:
            for mc in merge_cols:
                if mc < df.shape[1] and not pd.isna(df.iloc[data_row_idx, mc]):
                    is_kpi_first = True
                    break

        data_col_idx = 0
        for ppt_c in range(start_col, ppt_cols):
            if data_col_idx >= df.shape[1]:
                break
            if ppt_c in skip_cols:
                data_col_idx += 1
                continue
            value = df.iloc[data_row_idx, data_col_idx]
            if ppt_c in merge_cols and pd.isna(value):
                data_col_idx += 1
                continue
            cell = table.cell(ppt_r, ppt_c)
            text_value = format_value(value, data_col_idx, col_formats)

            # Determine numeric sign from raw value for reliable +/- detection.
            # Handle both plain numbers (-0.05) and already-formatted % strings ("-5%").
            try:
                s = str(value).replace(",", "").strip()
                if s.endswith("%"):
                    raw_num = float(s[:-1]) / 100
                else:
                    raw_num = float(s)
            except (ValueError, TypeError):
                raw_num = None

            is_negative = raw_num is not None and raw_num < 0
            is_positive = raw_num is not None and raw_num > 0

            # If this column requires pct-only colouring, skip non-percentage values
            is_pct_value = "%" in text_value
            skip_color = (ppt_c in pct_color_only_cols and not is_pct_value) \
                         or not row_color_allowed

            color = None
            if not skip_color and raw_num is not None:
                if color_threshold is not None:
                    # Threshold mode: >= threshold → green, < threshold → red
                    color = GREEN if raw_num >= color_threshold else RED
                else:
                    if (negative_red or ppt_c in neg_red_cols) and is_negative:
                        color = RED
                    elif (positive_green or ppt_c in pos_green_cols) and is_positive:
                        color = GREEN

            is_bold = None
            if ppt_c in always_bold_cols:
                is_bold = True
            elif data_row_idx in bold_data_rows:
                is_bold = True
            elif bold_first_in_merge:
                is_bold = is_kpi_first

            set_cell_text(cell, text_value, font_name=font_name, font_size=font_size,
                          bold=is_bold, font_color=color)
            data_col_idx += 1

        data_row_idx += 1

    print(f"  Updated table (index {table_idx}): "
          f"wrote {data_row_idx} data rows x {data_col_idx} cols "
          f"into a {ppt_rows}x{ppt_cols} PPT table.")

    apply_conditional_formatting(table, cfg, start_row, ppt_rows, ppt_cols)


def get_table_configs(config: dict, slide_num: int, sheet_count: int) -> list[dict]:
    """
    Build a list of per-table configs for a slide.
    If the config has a "tables" list, use that (one entry per table).
    Otherwise, wrap the flat config as a single-table list.
    """
    user_cfg = config.get(str(slide_num), {})

    if "tables" in user_cfg:
        table_list = user_cfg["tables"]
        return [{**DEFAULTS, **t_cfg} for t_cfg in table_list]

    merged = {**DEFAULTS, **user_cfg}
    return [merged]


def _load_full_config(config_file: Path = None) -> dict:
    """Load the raw slide_config.json including underscore keys."""
    cf = config_file or CONFIG_FILE
    if not cf.exists():
        return {}
    with open(cf, "r", encoding="utf-8") as fh:
        return json.load(fh)


def _save_full_config(raw: dict, config_file: Path = None):
    """Write the full config back to slide_config.json."""
    cf = config_file or CONFIG_FILE
    with open(cf, "w", encoding="utf-8") as fh:
        json.dump(raw, fh, indent=4, ensure_ascii=False)
    print(f"  Saved config to {cf.name}")


def _rename_data_file(folder: Path, old_num: int, new_num: int, dry_run: bool) -> str | None:
    """Rename a data file from one slide number to another. Returns description or None."""
    for ext in (".csv", ".xlsx"):
        for prefix in ("Slide", "slide"):
            old_path = folder / f"{prefix} {old_num}{ext}"
            if old_path.exists():
                new_path = folder / f"{prefix} {new_num}{ext}"
                if dry_run:
                    return f"  {old_path.name} -> {new_path.name}"
                old_path.rename(new_path)
                return f"  Renamed: {old_path.name} -> {new_path.name}"
    return None


def handle_insert(position: int, dry_run: bool = False, folder: Path = None, config_file: Path = None):
    """A new slide was inserted at `position`. Shift all configs and files >= position up by 1."""
    folder = folder or SCRIPT_DIR
    config_file = config_file or CONFIG_FILE
    raw = _load_full_config(config_file=config_file)
    data_files = find_data_files(folder)

    slide_keys = sorted(
        [int(k) for k in raw if not k.startswith("_") and k.isdigit()],
        reverse=True,
    )
    file_nums = sorted(data_files.keys(), reverse=True)

    action = "Would" if dry_run else "Will"
    print(f"\n{'[DRY RUN] ' if dry_run else ''}INSERT slide at position {position}")
    print(f"  {action} shift all slide configs and data files at position >= {position} up by 1.\n")

    changes = False

    print("Config changes:")
    for num in slide_keys:
        if num >= position:
            new_num = num + 1
            print(f"  Slide {num} -> Slide {new_num}")
            if not dry_run:
                raw[str(new_num)] = raw.pop(str(num))
            changes = True
    if not changes:
        print("  (none)")

    if not dry_run and changes:
        _save_full_config(raw, config_file=config_file)

    print("\nFile renames:")
    file_changes = False
    for num in file_nums:
        if num >= position:
            msg = _rename_data_file(folder, num, num + 1, dry_run)
            if msg:
                print(msg)
                file_changes = True
    if not file_changes:
        print("  (none)")

    print("\nDone." if not dry_run else "\n[DRY RUN] No changes made.")


def handle_remove(position: int, dry_run: bool = False, folder: Path = None, config_file: Path = None):
    """Slide at `position` was removed. Delete its config/file, shift everything above down by 1."""
    folder = folder or SCRIPT_DIR
    config_file = config_file or CONFIG_FILE
    raw = _load_full_config(config_file=config_file)
    data_files = find_data_files(folder)

    slide_keys = sorted(
        [int(k) for k in raw if not k.startswith("_") and k.isdigit()]
    )
    file_nums = sorted(data_files.keys())

    action = "Would" if dry_run else "Will"
    print(f"\n{'[DRY RUN] ' if dry_run else ''}REMOVE slide at position {position}")
    print(f"  {action} delete config/file for slide {position}, "
          f"then shift everything above it down by 1.\n")

    if str(position) in raw:
        print(f"Config: delete slide {position}")
        if not dry_run:
            del raw[str(position)]
    else:
        print(f"Config: slide {position} has no config entry (nothing to delete)")

    print("\nConfig shifts:")
    changes = False
    for num in slide_keys:
        if num > position:
            new_num = num - 1
            print(f"  Slide {num} -> Slide {new_num}")
            if not dry_run:
                raw[str(new_num)] = raw.pop(str(num))
            changes = True
    if not changes:
        print("  (none)")

    if not dry_run:
        _save_full_config(raw, config_file=config_file)

    print("\nFile changes:")
    if position in data_files:
        fpath = data_files[position]
        if dry_run:
            print(f"  Would delete: {fpath.name}")
        else:
            fpath.unlink()
            print(f"  Deleted: {fpath.name}")

    file_changes = False
    for num in file_nums:
        if num > position:
            msg = _rename_data_file(folder, num, num - 1, dry_run)
            if msg:
                print(msg)
                file_changes = True
    if not file_changes and position not in data_files:
        print("  (none)")

    print("\nDone." if not dry_run else "\n[DRY RUN] No changes made.")


def run_update(folder: Path, target_name: str = None, only_slides: set = None,
               mode: str = "individual", config_file: Path = None):
    """Core update logic: read data files, update PPT tables, save output.

    Args:
        folder:       Folder containing the PPT and data files.
        target_name:  PPT filename to update (auto-detected if None).
        only_slides:  If set, only process these slide numbers (individual mode).
        mode:         "individual" — use per-slide files (Slide N.xlsx / .csv)
                      "master"     — use GTM_GAI_Tabls.xlsx with sheets named
                                     "slide N" or "slide N.M"
        config_file:  Optional path to slide_config.json (defaults to SCRIPT_DIR/slide_config.json).
    """
    pptx_path = find_pptx(folder, target_name)
    print(f"PowerPoint: {pptx_path.name}")
    print(f"Data source mode: {mode}\n")

    config = load_config(config_file=config_file)
    prs = Presentation(str(pptx_path))
    total_slides = len(prs.slides)
    print(f"Presentation has {total_slides} slides.\n")

    # ── Build slide_sources: dict[slide_num -> list[sheet_names]] ──────────
    # Each entry is a list of sheet names (one per table on that slide).
    # For individual mode the "sheet name" is the actual sheet in the per-slide file.
    # For master mode it's the sheet name within GTM_GAI_Tabls.xlsx.

    if mode == "master":
        master_path = folder / MASTER_FILE_NAME
        slide_sources = parse_master_file(master_path)   # {slide_num: [sheet, ...]}
        # slide_sources values are (sheet_name, master_path) pairs below
        source_file: dict[int, Path] = {sn: master_path for sn in slide_sources}
        if not slide_sources:
            sys.exit(f"ERROR: No slide sheets found in {MASTER_FILE_NAME}.")
        print(f"Master file: {MASTER_FILE_NAME}")
        print(f"Found slide sheets: {sorted(slide_sources.keys())}\n")
    else:
        # individual mode
        data_files = find_data_files(folder)
        if only_slides:
            data_files = {k: v for k, v in data_files.items() if k in only_slides}
        if not data_files:
            sys.exit("No data files found (expected files like 'slide 1.csv', 'slide 15.xlsx').")
        print(f"Found data files for slides: {sorted(data_files.keys())}\n")
        # Convert to same structure as master: {slide_num: [sheet_names]}
        slide_sources = {}
        source_file = {}
        for sn, fp in data_files.items():
            slide_sources[sn] = get_sheet_names(fp)
            source_file[sn] = fp

    # ── Process each slide ─────────────────────────────────────────────────
    for slide_num in sorted(slide_sources.keys()):
        file_path  = source_file[slide_num]
        sheet_names = slide_sources[slide_num]
        print(f"--- Slide {slide_num} ({file_path.name}) ---")

        if slide_num < 1 or slide_num > total_slides:
            print(f"  WARNING: Slide {slide_num} is out of range (1-{total_slides}). Skipping.")
            continue

        slide      = prs.slides[slide_num - 1]
        table_cfgs = get_table_configs(config, slide_num, len(sheet_names))

        if len(sheet_names) > 1:
            print(f"  Found {len(sheet_names)} sheets: {sheet_names}")

        if len(sheet_names) > 1 and len(table_cfgs) == 1:
            print(f"  WARNING: {len(sheet_names)} sheets but config only has 1 table entry. "
                  f"Add a \"tables\" list in the config for slide {slide_num}.")
            print(f"  Processing only the first sheet with current config.\n")
            sheet_names = [sheet_names[0]]

        for i, sheet in enumerate(sheet_names):
            cfg = table_cfgs[i] if i < len(table_cfgs) else {**DEFAULTS}
            sheet_arg = None if file_path.suffix.lower() == ".csv" else sheet
            df = read_data(file_path, cfg, sheet_name=sheet_arg)

            label = f" (sheet: {sheet})" if len(sheet_names) > 1 else ""
            print(f"  [{sheet}] Data shape: {df.shape[0]} rows x {df.shape[1]} cols{label}")
            update_slide_table(slide, slide_num, df, cfg)

        print()

    output_path = pptx_path.with_name(
        pptx_path.stem + "_updated" + pptx_path.suffix
    )
    prs.save(str(output_path))
    print(f"Saved updated presentation to: {output_path.name}")
    return output_path


def main():
    args = sys.argv[1:]
    dry_run = "--dry-run" in args
    if dry_run:
        args.remove("--dry-run")

    if "--insert" in args:
        idx = args.index("--insert")
        if idx + 1 >= len(args):
            sys.exit("ERROR: --insert requires a slide number. Usage: py update_ppt.py --insert 13")
        position = int(args[idx + 1])
        handle_insert(position, dry_run=dry_run)
        return

    if "--remove" in args:
        idx = args.index("--remove")
        if idx + 1 >= len(args):
            sys.exit("ERROR: --remove requires a slide number. Usage: py update_ppt.py --remove 14")
        position = int(args[idx + 1])
        handle_remove(position, dry_run=dry_run)
        return

    # --mode master | individual  (default: individual)
    mode = "individual"
    if "--mode" in args:
        idx = args.index("--mode")
        if idx + 1 >= len(args):
            sys.exit("ERROR: --mode requires a value: master or individual")
        mode = args[idx + 1].lower()
        if mode not in ("master", "individual"):
            sys.exit(f"ERROR: --mode must be 'master' or 'individual', got '{mode}'")
        args = args[:idx] + args[idx + 2:]

    target_name = args[0] if args else None
    run_update(SCRIPT_DIR, target_name, mode=mode)


if __name__ == "__main__":
    main()
